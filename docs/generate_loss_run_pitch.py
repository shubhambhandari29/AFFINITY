from pathlib import Path
import sys

sys.path.insert(0, "/private/tmp/affinity_pptx_deps")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).parent.parent / "Loss_Run_Process_Understanding_and_Proposal.pptx"

NAVY = RGBColor(47, 47, 47)
# EXL-inspired orange-led presentation palette.
ORANGE = RGBColor(255, 75, 10)
BLUE = ORANGE
TEAL = RGBColor(221, 58, 5)
GOLD = RGBColor(255, 137, 54)
GREEN = RGBColor(54, 137, 89)
RED = RGBColor(187, 73, 73)
LIGHT_BLUE = RGBColor(255, 239, 231)
LIGHT_TEAL = RGBColor(255, 229, 217)
LIGHT_GOLD = RGBColor(255, 243, 233)
LIGHT_GREY = RGBColor(244, 246, 248)
MID_GREY = RGBColor(111, 122, 134)
DARK = RGBColor(37, 45, 54)
WHITE = RGBColor(255, 255, 255)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=20, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(slide, text, kicker=None):
    if kicker:
        add_text(slide, kicker.upper(), 0.65, 0.28, 4.5, 0.25, 9, TEAL, True)
    add_text(slide, text, 0.65, 0.56 if kicker else 0.38, 12.0, 0.55, 26, NAVY, True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.18), Inches(1.0), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()


def footer(slide, number, note="Discovery summary • Based on supplied notebook and SQL view definition"):
    add_text(slide, note, 0.65, 7.12, 10.8, 0.18, 8, MID_GREY)
    add_text(slide, f"{number:02d}", 12.05, 7.08, 0.55, 0.22, 9, NAVY, True, PP_ALIGN.RIGHT)


def rounded_box(slide, text, x, y, w, h, fill, line=None, color=DARK,
                size=16, bold=True, subtext=None):
    if isinstance(bold, str) and subtext is None:
        subtext = bold
        bold = True
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1.2)
    tf = shape.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.name = "Aptos"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    if subtext:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(5)
        r2 = p2.add_run(); r2.text = subtext; r2.font.name = "Aptos"; r2.font.size = Pt(max(9, size-5)); r2.font.color.rgb = color
    return shape


def arrow(slide, x, y, w=0.55, h=0.32, direction="right", color=GOLD):
    kind = {
        "right": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        "down": MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
    }[direction]
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    return shp


def pill(slide, text, x, y, w, fill, color=WHITE):
    return rounded_box(slide, text, x, y, w, 0.35, fill, fill, color, 10, True)


def bullet_list(slide, items, x, y, w, h, size=15, color=DARK, bullet_color=TEAL):
    # Native editable bullet-like list with colored markers.
    row_h = h / max(1, len(items))
    for i, item in enumerate(items):
        cy = y + i * row_h
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(cy + 0.11), Inches(0.10), Inches(0.10))
        dot.fill.solid(); dot.fill.fore_color.rgb = bullet_color; dot.line.fill.background()
        add_text(slide, item, x + 0.22, cy, w - 0.22, row_h, size, color)


# 1 — Cover
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide, WHITE)
accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.26), Inches(7.5))
accent.fill.solid(); accent.fill.fore_color.rgb = ORANGE; accent.line.fill.background()
corner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(10.70), Inches(0), Inches(2.63), Inches(0.34))
corner.fill.solid(); corner.fill.fore_color.rgb = ORANGE; corner.line.fill.background()
add_text(slide, "SPECIAL ACCOUNTS", 0.82, 0.72, 4.5, 0.3, 11, ORANGE, True)
add_text(slide, "Loss Run Generation", 0.82, 1.48, 9.8, 0.72, 36, NAVY, True)
add_text(slide, "Current-state understanding and application integration proposal", 0.82, 2.35, 10.8, 0.55, 20, DARK)
line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.84), Inches(3.20), Inches(1.35), Inches(0.07))
line.fill.solid(); line.fill.fore_color.rgb = ORANGE; line.line.fill.background()
add_text(slide, "Purpose", 0.82, 3.65, 1.2, 0.30, 14, ORANGE, True)
add_text(slide, "Document what the existing loss-run process does today and present a possible path to make it available through the current application.", 0.82, 4.08, 10.75, 1.10, 22, NAVY, True)
add_text(slide, "Prepared from SpecialAccountsToExcel.py and the supplied SAC_Loss_Run view definition", 0.84, 6.72, 11.3, 0.25, 9, MID_GREY)


# 2 — Executive understanding
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
title(slide, "Our understanding in one view", "Executive summary")
rounded_box(slide, "Database view", 0.75, 2.05, 2.25, 1.15, LIGHT_BLUE, BLUE, NAVY, 17,
            "Prepares report-ready claim and financial data")
arrow(slide, 3.14, 2.46)
rounded_box(slide, "Databricks notebook", 3.83, 2.05, 2.45, 1.15, LIGHT_TEAL, TEAL, NAVY, 17,
            "Formats and separates each customer's records")
arrow(slide, 6.42, 2.46)
rounded_box(slide, "Excel template", 7.12, 2.05, 2.20, 1.15, LIGHT_GOLD, GOLD, NAVY, 17,
            "Supplies cover, worksheets, pivots and chart")
arrow(slide, 9.46, 2.46)
rounded_box(slide, "Customer report", 10.16, 2.05, 2.30, 1.15, RGBColor(232, 244, 235), GREEN, NAVY, 17,
            "One editable .xlsx workbook per customer")
add_text(slide, "Core observation", 0.78, 4.05, 2.1, 0.28, 13, TEAL, True)
add_text(slide, "The database prepares the business data; Python mainly turns that data into the required Excel report.", 0.78, 4.43, 11.5, 0.65, 24, NAVY, True)
pill(slide, "CONFIRMED FROM CODE", 0.78, 5.62, 2.05, TEAL)
add_text(slide, "The current process is batch-oriented and runs outside the existing application.", 3.08, 5.59, 8.8, 0.42, 16, DARK)
footer(slide, 2)


# 3 — Source systems
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
title(slide, "Where the view lives—and where its data comes from", "Current state")
rounded_box(slide, "Same SQL Server / Managed Instance", 2.35, 1.40, 8.65, 0.58, NAVY, NAVY, WHITE, 17)
rounded_box(slide, "CLM_LakeHouse", 0.72, 2.45, 5.15, 3.20, LIGHT_BLUE, ORANGE, NAVY, 21,
            "Claims source database")
add_text(slide, "Underlying source tables", 1.14, 3.35, 4.30, 0.28, 14, ORANGE, True, PP_ALIGN.CENTER)
add_text(slide, "Claims • Policies • Financials\nClaimants • Vehicles • Litigation", 1.10, 3.82, 4.38, 0.78, 15, DARK, False, PP_ALIGN.CENTER)
add_text(slide, "The view reads these tables", 1.12, 4.90, 4.30, 0.28, 11, MID_GREY, True, PP_ALIGN.CENTER)
rounded_box(slide, "CLMAA_SpecialAccounts", 7.40, 2.45, 5.20, 3.20, LIGHT_TEAL, TEAL, NAVY, 21,
            "Special Accounts database")
rounded_box(slide, "tblAcctSpecial", 7.82, 3.47, 1.85, 1.10, WHITE, GOLD, NAVY, 14,
            "Eligible customers")
rounded_box(slide, "SAC_Loss_Run", 10.05, 3.47, 2.05, 1.10, ORANGE, ORANGE, WHITE, 15,
            "Reporting view lives here")
arrow(slide, 9.70, 3.86, 0.28, 0.25, "right", GOLD)
arrow(slide, 6.05, 3.72, 1.05, 0.48, "right", ORANGE)
add_text(slide, "cross-database read", 5.80, 4.28, 1.55, 0.28, 10, MID_GREY, True, PP_ALIGN.CENTER)
pill(slide, "CONFIRMED", 0.75, 6.22, 1.25, GREEN)
add_text(slide, "The view is stored in CLMAA_SpecialAccounts, but its SQL reads claims data from CLM_LakeHouse.", 2.17, 6.16, 10.0, 0.48, 15, DARK)
footer(slide, 3)


# 4 — Queries
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
title(slide, "What the notebook queries—and how each result is used", "Current state")
pill(slide, "QUERY 1", 0.72, 1.48, 1.08, ORANGE)
rounded_box(slide, "Eligible-customer query", 0.72, 1.96, 5.83, 0.72, LIGHT_BLUE, ORANGE, NAVY, 19)
add_text(slide, "Source", 0.98, 2.91, 0.80, 0.25, 11, ORANGE, True)
add_text(slide, "CLMAA_SpecialAccounts.dbo.tblAcctSpecial", 1.82, 2.88, 4.35, 0.32, 13, NAVY, True)
add_text(slide, "Filters", 0.98, 3.42, 0.80, 0.25, 11, ORANGE, True)
add_text(slide, "Active account • frequency present • frequency is not ‘Not Needed’", 1.82, 3.39, 4.45, 0.55, 13, DARK)
add_text(slide, "Returns", 0.98, 4.17, 0.80, 0.25, 11, ORANGE, True)
add_text(slide, "Customer number, customer name and loss-run frequency", 1.82, 4.14, 4.35, 0.52, 13, DARK)
add_text(slide, "Used for", 0.98, 4.88, 0.80, 0.25, 11, ORANGE, True)
add_text(slide, "Builds the list of customers processed by the notebook", 1.82, 4.85, 4.35, 0.52, 13, DARK)

pill(slide, "QUERY 2", 6.82, 1.48, 1.08, TEAL)
rounded_box(slide, "Loss-run data query", 6.82, 1.96, 5.80, 0.72, LIGHT_TEAL, TEAL, NAVY, 19)
add_text(slide, "Source", 7.08, 2.91, 0.80, 0.25, 11, TEAL, True)
add_text(slide, "CLMAA_SpecialAccounts.dbo.SAC_Loss_Run", 7.92, 2.88, 4.30, 0.32, 13, NAVY, True)
add_text(slide, "Statement", 7.08, 3.42, 0.80, 0.25, 11, TEAL, True)
add_text(slide, "SELECT * FROM SAC_Loss_Run", 7.92, 3.39, 4.30, 0.32, 13, DARK, True)
add_text(slide, "Returns", 7.08, 4.17, 0.80, 0.25, 11, TEAL, True)
add_text(slide, "Prepared policy, claim, feature, financial and report fields", 7.92, 4.14, 4.30, 0.52, 13, DARK)
add_text(slide, "Used for", 7.08, 4.88, 0.80, 0.25, 11, TEAL, True)
add_text(slide, "Filtered in Python by customer, then written into Excel", 7.92, 4.85, 4.30, 0.52, 13, DARK)
rounded_box(slide, "Important: the notebook loads the complete view result; the proposed UI flow could query only the selected customer.", 1.18, 5.78, 10.98, 0.70, LIGHT_GOLD, GOLD, NAVY, 15)
footer(slide, 4)


# 5 — View responsibility
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
title(slide, "What the SAC_Loss_Run view prepares", "Current state")
items = [
    ("Selects current records", "Current claim, feature and policy versions", BLUE),
    ("Combines business data", "Claims, policies, claimants, vehicles, adjusters and litigation", TEAL),
    ("Calculates financials", "Reserves, payments, ALAE, recoveries and incurred totals", GOLD),
    ("Applies report rules", "Open/closed claims, outstanding reserves and eligible accounts", GREEN),
]
for i, (head, body, color) in enumerate(items):
    y = 1.55 + i*1.18
    rounded_box(slide, str(i+1), 0.80, y, 0.60, 0.60, color, color, WHITE, 17)
    add_text(slide, head, 1.62, y-0.02, 3.2, 0.30, 17, NAVY, True)
    add_text(slide, body, 4.72, y-0.02, 7.45, 0.55, 14, DARK)
rounded_box(slide, "Resulting level", 0.82, 6.24, 1.75, 0.44, NAVY, NAVY, WHITE, 11)
add_text(slide, "Primarily one row per claim feature / exposure—not necessarily one row per claim.", 2.82, 6.20, 9.25, 0.42, 15, DARK)
footer(slide, 5)


# 6 — Workbook process/output
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
title(slide, "How the notebook turns view data into a workbook", "Current state")
workbook_steps = [
    ("1", "Customer selection", "Keeps only rows belonging to the current customer."),
    ("2", "Data formatting", "Preserves identifiers and formats exposure values."),
    ("3", "Report separation", "Splits normal claims from Record Only incidents."),
    ("4", "Template population", "Writes data into Claims Data and Record Only tables."),
    ("5", "Workbook updates", "Adds cover details and sets pivots to refresh on open."),
    ("6", "File creation", "Saves one dated .xlsx workbook for the customer."),
]
for i, (num, head, body) in enumerate(workbook_steps):
    row, col = divmod(i, 3)
    x, y = 0.72 + col*4.20, 1.52 + row*2.05
    rounded_box(slide, num, x, y, 0.52, 0.52, ORANGE, ORANGE, WHITE, 15)
    add_text(slide, head, x+0.70, y-0.01, 3.05, 0.30, 16, NAVY, True)
    add_text(slide, body, x+0.70, y+0.43, 3.10, 0.72, 12, DARK)
rounded_box(slide, "OUTPUT", 0.75, 5.75, 1.05, 0.42, NAVY, NAVY, WHITE, 11)
add_text(slide, "One Microsoft Excel workbook per qualifying customer", 2.02, 5.73, 4.55, 0.35, 15, NAVY, True)
add_text(slide, "Cover Page • Claims Data • Record Only • Summary by Policy Year • Chart", 6.55, 5.73, 5.90, 0.46, 13, DARK)
add_text(slide, "Filename: <Customer Name>_<YYYY_MM_DD>.xlsx", 2.02, 6.35, 6.3, 0.28, 11, MID_GREY)
footer(slide, 6)


# 7 — Proposed future state
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
title(slide, "How this could work within the existing application", "Proposed")
flow = [
    ("1", "User request", "User selects a customer and clicks Generate Loss Run."),
    ("2", "API validation", "FastAPI confirms authentication and the customer request."),
    ("3", "View query", "Application requests only that customer's loss-run rows."),
    ("4", "Excel creation", "Service applies current formatting and fills the template."),
    ("5", "File response", "The completed .xlsx workbook is returned to the user."),
]
for i, (num, head, body) in enumerate(flow):
    x = 0.55 + i*2.55
    rounded_box(slide, num, x+0.72, 1.48, 0.48, 0.48, ORANGE, ORANGE, WHITE, 14)
    rounded_box(slide, head, x, 2.08, 2.05, 0.66, NAVY if i in (0, 4) else ORANGE, NAVY if i in (0, 4) else ORANGE, WHITE, 14)
    add_text(slide, body, x, 2.94, 2.05, 0.90, 11, DARK, False, PP_ALIGN.CENTER)
    if i < len(flow)-1:
        arrow(slide, x+2.12, 2.29, 0.35, 0.25, "right", GOLD)
rounded_box(slide, "What stays", 0.78, 4.47, 1.35, 0.42, TEAL, TEAL, WHITE, 11)
add_text(slide, "SAC_Loss_Run view • business calculations • Excel layout and report rules", 2.35, 4.44, 9.70, 0.42, 14, DARK)
rounded_box(slide, "What changes", 0.78, 5.23, 1.35, 0.42, GOLD, GOLD, WHITE, 11)
add_text(slide, "User initiates the report • application queries one customer • workbook is delivered through the UI", 2.35, 5.20, 9.70, 0.50, 14, DARK)
pill(slide, "PROPOSAL — NOT YET IMPLEMENTED", 0.78, 6.27, 2.72, BLUE)
footer(slide, 7)


# 8 — considerations
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
title(slide, "Key considerations before delivery", "Discovery")
cards = [
    ("Template", "Confirm the approved workbook and required visual fidelity.", BLUE),
    ("Access", "Confirm application permission to read the reporting view.", TEAL),
    ("Performance", "Measure generation time for customers with large claim volumes.", GOLD),
    ("Security", "Confirm which users may generate which customer reports.", GREEN),
    ("Delivery", "Decide immediate download, storage, email—or a later phase.", NAVY),
    ("Validation", "Compare new output with a known-good current report.", RED),
]
for i, (head, body, color) in enumerate(cards):
    row, col = divmod(i, 3)
    x, y = 0.73 + col*4.18, 1.58 + row*2.18
    rounded_box(slide, head, x, y, 3.72, 0.50, color, color, WHITE, 14)
    rounded_box(slide, body, x, y+0.56, 3.72, 1.20, LIGHT_GREY, RGBColor(219, 224, 229), DARK, 13, False)
rounded_box(slide, "Recommended first scope", 0.75, 6.08, 2.15, 0.46, NAVY, NAVY, WHITE, 11)
add_text(slide, "Generate one selected customer's workbook and return it to the authenticated user.", 3.15, 6.08, 9.05, 0.42, 15, DARK)
footer(slide, 8)


# 9 — confirmation register
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
title(slide, "What is confirmed—and what remains open", "Client discussion")
rounded_box(slide, "CONFIRMED FROM SUPPLIED MATERIAL", 0.72, 1.48, 5.90, 0.54, GREEN, GREEN, WHITE, 13)
confirmed = [
    "Two databases on the same server / managed instance",
    "SAC_Loss_Run supplies the prepared report data",
    "Existing workbook template is populated per customer",
    "Normal claims and Record Only incidents are separated",
    "Output is an editable Microsoft Excel workbook",
]
bullet_list(slide, confirmed, 0.90, 2.20, 5.45, 3.30, 14, DARK, GREEN)
rounded_box(slide, "TO CONFIRM WITH CLIENT", 6.78, 1.48, 5.82, 0.54, GOLD, GOLD, WHITE, 13)
open_items = [
    "One-customer download, batch generation, or both",
    "Immediate download versus stored or emailed report",
    "Approved template and exact output-matching expectation",
    "Business rules and date-range expectations",
    "Expected user permissions and report volumes",
]
bullet_list(slide, open_items, 6.98, 2.20, 5.35, 3.30, 14, DARK, GOLD)
rounded_box(slide, "Proposed next step", 3.72, 6.08, 2.05, 0.48, NAVY, NAVY, WHITE, 12)
add_text(slide, "Validate one representative customer and agree the target user experience.", 6.02, 6.08, 5.72, 0.45, 15, NAVY, True)
footer(slide, 9)


# 10 — technical appendix
slide = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(slide)
title(slide, "Technical source reference", "Appendix")
add_text(slide, "Claims and reference sources — CLM_LakeHouse", 0.75, 1.45, 5.8, 0.32, 16, BLUE, True)
left = [
    "DW_CLMS_CLM_DIMNSN — claims",
    "DW_CLMS_FTR_DIMNSN — claim features",
    "DW_FINCL_LOSS_FTR_SS_F — financial feature totals",
    "DW_FINCL_LOSS_TXN_SS_F — financial transactions",
    "DW_POL_DIMNSN — policies",
    "DW_CLMS_CLMNT_DIMNSN — claimants",
]
right = [
    "DW_CLMS_DRVR_DIMNSN / CLM_CNTCT_DIMNSN — drivers",
    "DW_CLMS_VEH_DIMNSN — vehicles",
    "DW_CLMS_LGL_FACT / LGL_DIMNSN — litigation",
    "DW_REF_MAJ_PERIL — peril reference",
    "DM_RETNTN_PRDCNG_AGNT_HR_D — producing agent",
]
bullet_list(slide, left, 0.80, 1.95, 5.95, 3.55, 12, DARK, BLUE)
bullet_list(slide, right, 6.67, 1.95, 5.90, 3.55, 12, DARK, BLUE)
rounded_box(slide, "CLMAA_SpecialAccounts", 0.78, 5.66, 2.45, 0.48, TEAL, TEAL, WHITE, 12)
add_text(slide, "tblAcctSpecial — eligible customer configuration   •   SAC_Loss_Run — reporting view", 3.48, 5.66, 8.72, 0.46, 14, DARK)
rounded_box(slide, "Notebook assets", 0.78, 6.35, 2.45, 0.48, GOLD, GOLD, WHITE, 12)
add_text(slide, "SACLossRunTemplate.xlsx   •   Databricks Unity Catalog volume   •   .xlsx output", 3.48, 6.35, 8.72, 0.46, 14, DARK)
footer(slide, 10, "Appendix • Technical object names from supplied materials")


prs.core_properties.title = "Special Accounts Loss Run Generation — Understanding and Proposal"
prs.core_properties.subject = "Current-state flow and potential FastAPI application integration"
prs.core_properties.author = "Affinity project team"
prs.core_properties.keywords = "loss run, Special Accounts, FastAPI, Databricks, Excel"
prs.save(OUT)
print(OUT)
